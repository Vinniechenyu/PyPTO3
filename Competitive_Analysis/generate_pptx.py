# -*- coding: utf-8 -*-
"""生成 PyPTO 开发者体验(DevUX)竞品分析 PPT。
所有内容严格取自项目内已核实文档（竞品资料收集审核稿、PyPTO3.0 产品功能规划、hw-native-sys 仓库功能洞察报告），
规划/推断内容均显式标注，不编造任何数据。
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- 设计令牌 ----------
PRIMARY = RGBColor(0x1F, 0x38, 0x64)   # 深蓝
ACCENT  = RGBColor(0x2E, 0x75, 0xB6)   # 蓝
LIGHT   = RGBColor(0xF2, 0xF4, 0xF8)   # 浅灰底
BAND    = RGBColor(0xE8, 0xEE, 0xF6)   # 表头浅蓝
GREY    = RGBColor(0x59, 0x59, 0x59)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
RED     = RGBColor(0xC0, 0x00, 0x00)   # 待核验/警示
GREEN   = RGBColor(0x2E, 0x7D, 0x32)
AMBER   = RGBColor(0xB7, 0x6E, 0x00)
FONT    = "Microsoft YaHei"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def set_font(run, size=14, bold=False, color=GREY, font=FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # 同时设置 latin 与 ea（CJK）字体，确保中文显示
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", font)

def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp

def add_text(slide, x, y, w, h, text, size=14, bold=False, color=GREY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT, line_spacing=1.05):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run(); r.text = text
    set_font(r, size, bold, color, font)
    return tb

def add_bullets(slide, x, y, w, h, items, size=14, gap=6, base_color=GREY):
    """items: list of (text, level, color, bold)"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    first = True
    for it in items:
        text, level, color, bold = (it + (base_color, False))[:4] if len(it) < 4 else it
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(gap)
        p.line_spacing = 1.05
        marker = "•  " if level == 0 else "–  "
        r = p.add_run(); r.text = marker + text
        set_font(r, size - level, bold, color)
    return tb

def header(slide, title, kicker=None):
    add_rect(slide, 0, 0, SW, Inches(1.05), PRIMARY)
    add_rect(slide, 0, Inches(1.05), SW, Pt(3), ACCENT)
    if kicker:
        add_text(slide, Inches(0.5), Inches(0.12), Inches(12), Inches(0.3),
                 kicker, size=11, bold=True, color=RGBColor(0xBD,0xD3,0xEE))
    add_text(slide, Inches(0.5), Inches(0.38) if kicker else Inches(0.28),
             Inches(12.3), Inches(0.6), title, size=24, bold=True, color=WHITE)
    # 页脚
    add_text(slide, Inches(0.5), SH - Inches(0.32), Inches(9), Inches(0.25),
             "PyPTO 开发者体验(DevUX)竞品分析 · 数据均来自项目内已核实资料", size=9, color=GREY)
    add_text(slide, SW - Inches(1.2), SH - Inches(0.32), Inches(0.8), Inches(0.25),
             "2026-08-11", size=9, color=GREY, align=PP_ALIGN.RIGHT)

def new_slide():
    return prs.slides.add_slide(BLANK)

def style_table(table, header_fill=BAND, header_color=PRIMARY, body_size=11,
                header_size=11, col_colors=None):
    # 表头
    for j, cell in enumerate(table.rows[0].cells):
        cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Pt(4); cell.margin_right = Pt(4)
        cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                set_font(r, header_size, True, header_color)
    # 表体
    for i in range(1, len(table.rows)):
        for j, cell in enumerate(table.rows[i].cells):
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Pt(4); cell.margin_right = Pt(4)
            cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    c = GREY
                    if col_colors and j in col_colors:
                        c = col_colors[j]
                    set_font(r, body_size, False, c)

def add_table(slide, x, y, w, h, headers, rows, col_widths=None,
              body_size=11, header_size=11):
    nrows = len(rows) + 1
    ncols = len(headers)
    gfx = slide.shapes.add_table(nrows, ncols, x, y, w, h)
    table = gfx.table
    # 关闭默认样式条带
    tbl = table._tbl
    # 设置列宽
    if col_widths:
        total = sum(col_widths)
        for j, cw in enumerate(col_widths):
            table.columns[j].width = Emu(int(w * cw / total))
    # 表头
    for j, htext in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = htext
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = "" if val is None else str(val)
    style_table(table, body_size=body_size, header_size=header_size)
    return table

def note_line(slide, x, y, w, text, color=AMBER):
    add_text(slide, x, y, w, Inches(0.3), text, size=10, color=color, bold=True)

# =========================================================
# Slide 1 — 封面
# =========================================================
s = new_slide()
add_rect(s, 0, 0, SW, SH, PRIMARY)
add_rect(s, 0, Inches(4.55), SW, Pt(3), ACCENT)
add_text(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.0),
         "PyPTO 开发者体验(DevUX)竞品分析", size=40, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(2.7), Inches(11.7), Inches(0.6),
         "基于项目内已核实资料的事实与大纲", size=20, color=RGBColor(0xBD,0xD3,0xEE))
add_text(s, Inches(0.8), Inches(4.8), Inches(11.7), Inches(1.6),
         "数据来源：PyPTO3_竞品资料收集审核稿(2026-08-05) · PyPTO3.0 Toolkit 产品功能规划 · "
         "hw-native-sys 仓库功能洞察报告。\n"
         "原则：所有数据/结论均来自上述已核实文档；规划中或推断内容已明确标注，不编造任何数字。",
         size=13, color=RGBColor(0xD6,0xE2,0xF1), line_spacing=1.25)
add_text(s, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.4),
         "Competitive_Analysis · 2026-08-11", size=12, color=RGBColor(0x9D,0xB8,0xD8))

# =========================================================
# Slide 2 — 分析目的与范围
# =========================================================
s = new_slide()
header(s, "分析目的与范围", "WHY DEVUX")
add_bullets(s, Inches(0.5), Inches(1.35), Inches(12.3), Inches(2.4), [
    ("目的：从“功能/技术对比”收敛到“开发者体验对比”——开发者在昇腾 NPU 上，从拿到任务到产出可信、可复现、可交付的结果，用 PyPTO 与竞品的体验差在哪里。", 0, GREY, False),
    ("北极星指标 TTTT（Time to Trusted Target）：从首次运行到“环境可复现 + 正确性门禁通过 + 性能/资源达标 + 证据完整”的可信产物所需中位时间。（来源：产品功能规划）", 0, GREY, False),
    ("三层分析结构（复用审核稿）：核心直接竞品 / 战略标杆 / 邻近参照。", 0, GREY, False),
], size=14, gap=10)
note_line(s, Inches(0.5), Inches(3.95), Inches(12.3), "范围说明（重要）")
add_bullets(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.6), [
    ("分析对象严格区分“当前可用 PyPTO”与“PyPTO 3.0 目标产品”；本文事实均来自已核实资料，规划项单独标注。", 0, GREY, False),
    ("竞品数据与判断来自审核稿中的官方文档 / 官方 GitHub 引用；待核验项以【待核验】标注，未做实测前不写结论。", 0, GREY, False),
    ("生态数字（stars/forks/issue）仅作生态成熟度辅助信号，不作为技术结论。", 0, GREY, False),
], size=13, gap=8)

# =========================================================
# Slide 3 — 竞品分层总览
# =========================================================
s = new_slide()
header(s, "竞品分层总览（10 家）", "LANDSCAPE")
rows = [
    ["核心直接竞品", "TileLang-Ascend", "Ascend 上写高性能 kernel 最直接对手；Developer/Expert 双模式、自动同步、自动 buffer reuse、wheel、课程体系"],
    ["核心直接竞品", "Triton-Ascend", "迁移成本最低、开发者心智复用最强（将成熟 Triton 范式带到 Ascend）"],
    ["核心直接竞品", "CANN Ascend C + MindStudio", "Ascend 原生算子开发与调优完整工具链，PyPTO 用户最现实的替代路径"],
    ["核心直接竞品", "Triton（上游）", "GPU kernel DSL 事实心智基准；Gluon 补足低层 tile-based SPMD 控制"],
    ["核心直接竞品", "TileLang（上游）", "多后端、开放、生态动能强（约 6.4k stars）；已用于 BitBLAS / AttentionEngine"],
    ["战略标杆", "cuTile Python + Nsight Compute", "Python Tile DSL + 编译器 + profiler + 生态闭环的产品化标杆"],
    ["战略标杆", "CUTLASS / CuTe DSL", "Tile/Block/ISA 分层理念最接近的 NVIDIA 参照"],
    ["邻近参照", "Apache TVM TensorIR", "调优搜索空间 / 实验数据库 / schedule 可追踪性参照"],
    ["邻近参照", "IREE", "模型→runtime 工程化、低开销运行时参照"],
    ["邻近参照", "vLLM / SGLang / TensorRT-LLM", "服务层 DevUX 参照（请求重放、容量、SLO 门禁）"],
]
add_table(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.6),
          ["层次", "竞品", "在 DevUX 上的最强对标点"], rows,
          col_widths=[1.6, 3.0, 7.7], body_size=11, header_size=12)

# =========================================================
# Slide 4 — PyPTO 己方基线
# =========================================================
s = new_slide()
header(s, "PyPTO 己方基线：事实、规划与软肋", "BASELINE")
add_text(s, Inches(0.5), Inches(1.25), Inches(6), Inches(0.3), "已核验事实（审核稿 §2.1）", size=13, bold=True, color=PRIMARY)
add_bullets(s, Inches(0.5), Inches(1.6), Inches(6.1), Inches(3.0), [
    ("面向 AI 加速器的高性能编程框架，PTO + Tile 编程模型为核心", 0),
    ("多级 IR：Tensor→Tile→Block→Execution Graph→PTO 虚拟指令", 0),
    ("支持 MPMD 执行调度", 0),
    ("分层抽象：算法(Tensor)/性能专家(Tile)/系统(Block+指令)", 0),
    ("公开实现主要面向华为 AI 处理器；许可证限定衍生软件用于华为 AI 处理器系统", 0),
    ("公开仓库(2026-08-05)：约 1384 commits、102 stars、82 forks、63 open issue、39 open PR", 0),
], size=11.5, gap=5)
add_text(s, Inches(6.9), Inches(1.25), Inches(6), Inches(0.3), "PyPTO 3.0 规划假设（规划中，非已交付）", size=13, bold=True, color=AMBER)
add_bullets(s, Inches(6.9), Inches(1.6), Inches(5.9), Inches(2.2), [
    ("Development Evidence Graph：关联 request→source→IR→ISA→runtime→tensor→metric", 0, AMBER, False),
    ("Compile Guardian / Pass Contract：各编译阶段验证不变量，优先暴露静默错误", 0, AMBER, False),
    ("Correctness Lab / Performance Lab / Inference Bundle / Service Builder", 0, AMBER, False),
    ("北极星指标 TTTT", 0, AMBER, False),
], size=11.5, gap=5)
add_text(s, Inches(6.9), Inches(3.95), Inches(6), Inches(0.3), "已识别软肋（洞察报告）", size=13, bold=True, color=RED)
add_bullets(s, Inches(6.9), Inches(4.3), Inches(5.9), Inches(2.4), [
    ("仓库命名混乱（simpler 实为 PTO Runtime）、文档分散", 0, RED, False),
    ("issue 集中在 pypto，反映“易用性是主要痛点”", 0, RED, False),
    ("Fork/Star 高，属内部工程组织而非社区驱动", 0, RED, False),
], size=11.5, gap=5)
note_line(s, Inches(0.5), Inches(4.75), Inches(6), "事实 vs 规划：左侧为已交付/已核验；右上为规划方向，右下为已观察短板。")

# =========================================================
# Slide 5-9 — 竞品详解
# =========================================================
def competitor_slide(title, kicker, why, judge, pending):
    s = new_slide()
    header(s, title, kicker)
    add_text(s, Inches(0.5), Inches(1.25), Inches(12), Inches(0.3), "为什么最直接 / 关键事实", size=13, bold=True, color=PRIMARY)
    add_bullets(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(2.2), why, size=12.5, gap=6)
    add_text(s, Inches(0.5), Inches(3.95), Inches(12), Inches(0.3), "初步判断", size=13, bold=True, color=ACCENT)
    add_bullets(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(1.6), judge, size=12.5, gap=6)
    note_line(s, Inches(0.5), Inches(5.95), Inches(12.3), "待核验（未实测，不做结论）: " + pending, color=RED)
    return s

competitor_slide(
    "核心竞品① TileLang-Ascend", "CORE 1/5",
    [("同为 Ascend NPU 上的 Python/Tile DSL，用户群、硬件环境与算子类型高度重叠", 0),
     ("已提供 Developer/Expert 双控制、自动同步、自动 buffer reuse、软件流水、Cube/Vector 分离、PTO backend、PyTorch/ACLGraph 集成", 0),
     ("官方示例含 Flash Attention、Sparse Attention、GEMM、归约、卷积、MoE dispatch/combine；明确测试 A2/A3", 0),
     ("MIT 许可证、wheel 安装、课程体系，对开发者采用有直接影响", 0)],
    [("是 PyPTO 在“Ascend 上写高性能 kernel”场景中最强的正面竞品", 0, GREY, False),
     ("PyPTO 若只讲 Tile DSL 和自动同步，差异不够明显；需突出多层抽象、MPMD 编排、PTO ISA 一致性，以及 source→runtime 可信证据链", 0, GREY, False)],
    "同硬件、同 shape、同精度下的性能、编译时间、首个 kernel 成功时间、调试效率"
)

competitor_slide(
    "核心竞品② Triton-Ascend", "CORE 2/5",
    [("将成熟的 Triton 编程范式带到 Ascend，直接争夺已有 Triton/PyTorch 开发者", 0),
     ("官方定位：由编译器自动完成内存分配、数据搬运、计算与流水并行，降低算子开发门槛", 0),
     ("当前公开信息显示兼容 CANN 8.5，并有 pip 安装路径；路线图含向更新 Triton 版本演进", 0)],
    [("最大优势不是 Ascend 特有的最强控制力，而是迁移成本与开发者心智复用", 0, GREY, False),
     ("PyPTO 可从“Ascend 原生语义深度、Tensor/Tile/Block 分层、跨核/跨卡编排、硬件约束可解释”建立区隔", 0, GREY, False)],
    "API 覆盖率、生产案例、复杂融合算子性能稳定性、调试工具与多卡能力"
)

competitor_slide(
    "核心竞品③ CANN Ascend C + MindStudio", "CORE 3/5",
    [("Ascend 原生算子开发与调优的现有完整工具链，也是 PyPTO 用户最现实的替代路径", 0),
     ("MindStudio Operator Tools 覆盖性能建模、项目生成、功能测试、异常检测、板上/仿真调试、性能采集", 0),
     ("MindStudio Insight 支持系统/算子/服务/内存调优，含指令流水、源码、负载、集群时间线视图", 0)],
    [("强项：官方性、硬件覆盖、诊断深度、既有流程", 0, GREY, False),
     ("弱项更可能是开发门槛、工具割裂、跨层因果解释成本，而非能力缺失", 0, GREY, False),
     ("定位为“能力强但开发路径重”的原生基线，不宜描述为“落后工具”", 0, GREY, False)],
    "最新 CANN/MindStudio 版本的实际安装、许可、可用硬件矩阵、典型工作流耗时"
)

competitor_slide(
    "核心竞品④⑤ Triton（上游） / TileLang（上游）", "CORE 4-5/5",
    [("Triton：Python GPU kernel DSL 事实心智基准；教程、PyTorch/Inductor 生态、大量现成 kernel 构成采用优势；新增 Gluon 提供更低层 tile-based SPMD 模型", 0),
     ("TileLang：基于 TVM/TIR 的高性能 kernel DSL，面向 GPU/CPU/加速器；支持显式内存层级、流水、硬件特性、JIT；约 6.4k stars、582 forks，已用于 BitBLAS/AttentionEngine", 0)],
    [("Triton+Gluon 正形成“高低层组合”，PyPTO 差异化窗口收窄；论证重点应从“比 Triton 多一层”升级为“在异构 NPU 上贯通语义/验证/执行/证据”", 0, GREY, False),
     ("TileLang 优势是开放、多后端、快速演进；PyPTO 优势应是 Ascend 原生语义、整栈协同、可验证性——“硬件专用深度 vs 多后端广度”是战略取舍", 0, GREY, False)],
    "上游版本在 Ascend 路径的覆盖与演进节奏；多后端一致性的真实成本"
)

competitor_slide(
    "战略标杆 cuTile + Nsight / CUTLASS-CuTe", "STRATEGIC",
    [("cuTile：数组/Tile 驱动的 Python 编程模型，边界可检查、无裸指针、JIT/AOT、TileIR 导出、autotune、JAX 集成；执行模型仅暴露 block 级并行", 0),
     ("Nsight Compute：kernel 级指标、API 调试、图形报告、baseline 比较、参数系列实验、可扩展分析规则", 0),
     ("CUTLASS 4.x 同时提供 C++ 模板与 Python DSL；CuTe DSL 暴露 layout/tensor/hardware atom/tiled operation；最新文档加入源码→PTX/SASS 关联、调试模式、内核内事件追踪", 0)],
    [("NVIDIA 正把“Python Tile DSL + 编译器 + profiler + 生态”产品化，与 PyPTO 3.0 方向高度同构——可借鉴安全数据模型、AOT、实验对比、工具链闭环", 0, GREY, False),
     ("CUTLASS/CuTe 是 PyPTO Tile/Block/ISA 分层理念最接近的 NVIDIA 参照；对比重点应是专家控制上限、调试证据、可组合库、生产集成", 0, GREY, False)],
    "cuTile 在 Ascend 类硬件上的可迁移性；Nsight 式闭环的构建成本"
)

# =========================================================
# Slide 10 — 邻近参照
# =========================================================
s = new_slide()
header(s, "邻近参照（服务层 / 编译层，仅在相关能力页出现）", "ADJACENT")
rows = [
    ["Apache TVM TensorIR + MetaSchedule", "TensorIR 提供 tensor function 表示/变换/schedule primitives；MetaSchedule 以真实硬件测量、进化搜索、成本模型、持久化数据库寻找更优 schedule，支持跨模型复用", "调优搜索空间、实验数据库、schedule 可追踪性；非日常开发入口的完全直接竞品"],
    ["IREE", "基于 MLIR 的端到端编译器与 runtime，从框架模型下降到统一 IR，支持 AOT 与多 CPU/GPU/边缘平台；工具支持逐 pass 运行/检查/执行/dump", "模型→runtime 工程化、部署配置、低开销运行时；kernel 手写体验上不与 PyPTO 完全同层"],
    ["vLLM / SGLang / TensorRT-LLM", "vLLM：OpenAI 兼容在线推理；SGLang：在线 benchmark 覆盖 TTFT/ITL/吞吐/并发，支持请求 dump/replay、crash dump/replay、profiling；TRT-LLM：KV cache/chunked prefill/并行/低精度/分离式服务", "PyPTO 3.0 的 Service Builder/Workload Lab/请求重放/容量规划/SLO 门禁；服务层参照，非核心 DSL 直接竞品"],
]
add_table(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(4.6),
          ["参照对象", "能力事实", "参照价值（定位）"], rows,
          col_widths=[2.6, 6.0, 3.7], body_size=11, header_size=12)
note_line(s, Inches(0.5), Inches(6.05), Inches(12.3), "这些项目是能力参照，不宜被描述为 PyPTO 核心 DSL 的直接竞品。")

# =========================================================
# Slide 11 — DevUX 评估主线：开发者旅程九阶段
# =========================================================
s = new_slide()
header(s, "DevUX 评估主线：开发者旅程九阶段", "DEV JOURNEY")
rows = [
    ["A 上手与环境", "安装/环境校验/无硬件可用", "pypto doctor、Compatibility Profile、Resource Preflight；无硬件模拟器"],
    ["B 编写与表达", "DSL 人机工程/IDE/模板", "意图式 DSL、语义化内存 API、Intent Preview、真实任务模板（规划）"],
    ["C 编译与可信", "静态期拦截静默错误", "Pass Contract / verifier / fail-loud（规划）"],
    ["D 运行与调试", "source↔runtime 关联/错误字典", "Evidence Graph / Provenance Explorer / Error Dictionary（规划）"],
    ["E 正确性", "多 oracle / 首个分歧", "Correctness Lab / Runtime Sentinel（规划）"],
    ["F 性能", "跨层因果 / 可信 A/B", "四层指标 / Experiment Board / Safe Autotune（规划）"],
    ["G 分布式", "collective 校验 / hang 定位", "跨 rank verifier / 通信-计算 overlap（规划）"],
    ["H 模型服务", "导入→服务 / 请求追踪", "Model Importer / Inference Bundle / Service Builder（部分早期）"],
    ["I 复现与知识", "Repro / 文档可执行", "Repro Bundle / Issue Router / Developer Portal（规划）"],
]
add_table(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.4),
          ["阶段", "关注点", "PyPTO 基线（规划项已标注）"], rows,
          col_widths=[2.3, 3.4, 6.6], body_size=11, header_size=12)
note_line(s, Inches(0.5), Inches(6.75), Inches(12.3), "标注“（规划）”者为 PyPTO 3.0 目标能力，非当前已交付。")

# =========================================================
# Slide 12 — DevUX 评估维度（12 维度 + 规划目标值）
# =========================================================
s = new_slide()
header(s, "DevUX 评估维度（12 项）与规划目标值", "METRICS")
rows = [
    ["1 上手成本", "安装/环境校验/无硬件可用", "首次 smoke test 成功率 ≥90%（规划目标）"],
    ["2 表达效率", "DSL 人机工程/IDE/模板", "首个复杂 kernel 代码量、返工次数"],
    ["3 编译可信", "静态期拦截静默错误", "高危静默错误捕获率 ≥80%（规划目标）"],
    ["4 调试可解释", "source↔IR↔runtime 关联", "首个异常自动定位率 ≥70%；中位定位时间降 60%（规划目标）"],
    ["5 正确性闭环", "多 oracle / 首个分歧", "分歧定位时间"],
    ["6 性能可归因", "跨层因果 / 可信 A/B", "回退归因到具体层级比例 ≥75%（规划目标）"],
    ["7 分布式 DevUX", "collective 校验 / hang 定位", "多卡问题定位时间"],
    ["8 服务化 DevUX", "导入→服务 / 请求追踪", "导入→首个正确 token ≤30min；服务首次响应 ≤10min（规划目标）"],
    ["9 复现与协作", "Repro / 文档可执行", "问题单完整证据比例 ≥85%；裸错误码占比 <5%（规划目标）"],
    ["10 文档与知识", "可执行文档 / 角色路径", "文档 CI 可执行通过率 100%（规划目标）"],
    ["11 生态与治理", "许可证/硬件锁定/社区", "stars/forks/issue 密度（辅助信号）"],
    ["12 信任叙事", "fail-loud / 可解释 / 可复现", "定性，结合 TTTT"],
]
add_table(s, Inches(0.5), Inches(1.25), Inches(12.3), Inches(5.7),
          ["维度", "关注点", "指标（规划目标值，非实测）"], rows,
          col_widths=[2.2, 3.6, 6.5], body_size=10.5, header_size=12)
note_line(s, Inches(0.5), Inches(6.98), Inches(12.3), "括号内“规划目标”来自产品功能规划 §9.2，是建议目标值，不代表当前已达成。")

# =========================================================
# Slide 13 — 定量评测场景
# =========================================================
s = new_slide()
header(s, "定量评测场景（映射到 DevUX 指标）", "BENCH")
rows = [
    ["首个 kernel", "环境准备→vector add/matmul 正确运行时间", "同硬件同 shape 对比"],
    ["复杂 kernel", "Flash Attn / paged attn 代码量、调优轮数、性能", "TileLang-Ascend / Triton-Ascend 同机对比"],
    ["错值定位", "注入 layout/shape/sync 错误，定位首个分歧点时间", "PyPTO vs 竞品"],
    ["性能回归", "对同 kernel 做一次可控退化，比较瓶颈定位时间", "跨层归因能力"],
    ["动态 shape", "多 shape 编译缓存、正确性与性能稳定性", "—"],
    ["多核/多卡", "通信/同步/hang 定位能力", "—"],
    ["模型闭环", "模型导入→首个正确 token→满足 TTFT/TPOT", "服务层参照 vLLM/SGLang"],
    ["可复现", "另一台兼容机同 bundle 重放结果", "脱敏与复现包"],
]
add_table(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.0),
          ["场景", "DevUX 指标", "待核验"], rows,
          col_widths=[2.2, 6.0, 4.1], body_size=11, header_size=12)
note_line(s, Inches(0.5), Inches(6.35), Inches(12.3), "以上场景均需在真机（A2/A3/A5）实测后才能得出对比结论，当前为评测设计，非结果。")

# =========================================================
# Slide 14 — 核心论点与差异化叙事
# =========================================================
s = new_slide()
header(s, "核心论点与差异化叙事（DevUX 版）", "THESIS")
add_bullets(s, Inches(0.5), Inches(1.35), Inches(12.3), Inches(4.6), [
    ("PyPTO 的真正竞争对象不是单一 Triton，而是“Ascend C+MindStudio 的原生深度、Triton/TileLang 的开发效率、NVIDIA 工具链的闭环体验”的组合。", 0),
    ("TileLang-Ascend 是当下最直接的产品竞争者；Triton-Ascend 是最危险的迁移入口竞争者。", 0),
    ("潜在护城河不是 Python 语法，而是 PTO ISA 之上的多层抽象 + MPMD 执行 + 跨层可信证据（“知道结果从哪来”）。", 0),
    ("多层抽象本身已不再独占（Triton Gluon、CuTe DSL 正补足低层控制）；“可验证、可追溯、可复现”应成为 3.0 的核心产品叙事。", 0),
    ("仅限华为 AI 处理器的许可证既强化生态归属，也削弱通用开发者采用与跨硬件迁移能力——应正面呈现这一取舍。", 0),
    ("服务化能力应作为 PyPTO 整栈闭环的延伸，而非与 vLLM/SGLang 正面对标全部生产 serving 能力。", 0),
], size=13, gap=9)
note_line(s, Inches(0.5), Inches(6.15), Inches(12.3), "以上观点来自审核稿 §5，均需用户审核后才能进入最终稿；本页未新增任何数据。", color=AMBER)

# =========================================================
# Slide 15 — 材料收集进度与待确认
# =========================================================
s = new_slide()
header(s, "材料收集进度与待确认", "STATUS")
add_text(s, Inches(0.5), Inches(1.2), Inches(6), Inches(0.3), "收集状态（✅已收集 / 🔶部分待补 / ⬜待收集）", size=13, bold=True, color=PRIMARY)
rows = [
    ["己方基线 / 竞品清单 / 维度 / 场景", "✅", "来自审核稿、规划、洞察报告"],
    ["竞品 DevUX 对标点（分层）", "✅", "审核稿 §3"],
    ["上手/性能/服务 部分细节", "🔶", "需抓取各官方 doc / Nsight / MindStudio"],
    ["编译可信 / 正确性 / 分布式 / 复现 DevUX", "⬜", "需逐家核实 verifier / oracle / 多卡"],
    ["竞品×阶段矩阵、DevUX 雷达图", "⬜", "依赖上述填充后生成"],
]
add_table(s, Inches(0.5), Inches(1.6), Inches(7.2), Inches(3.0),
          ["内容", "状态", "来源/说明"], rows,
          col_widths=[3.6, 1.0, 2.6], body_size=10.5, header_size=11)
add_text(s, Inches(8.0), Inches(1.2), Inches(4.8), Inches(0.3), "待确认（阻塞正式材料）", size=13, bold=True, color=RED)
add_bullets(s, Inches(8.0), Inches(1.6), Inches(4.8), Inches(4.4), [
    ("受众与用途：管理层/产品规划/研发架构/生态/对外？", 0, RED, False),
    ("分析对象：当前可用 PyPTO 还是 3.0 目标？", 0, RED, False),
    ("能否评价 CANN/MindStudio 内部产品、可否用实测数据？", 0, RED, False),
    ("是否有 A2/A3/A5 真机做 benchmark？", 0, RED, False),
    ("交付形态：PPT / Word / Markdown？", 0, RED, False),
], size=12, gap=8)

# =========================================================
# Slide 16 — 公开来源清单
# =========================================================
s = new_slide()
header(s, "公开来源清单（DevUX 相关）", "SOURCES")
add_text(s, Inches(0.5), Inches(1.2), Inches(6), Inches(0.3), "己方 / 核心竞品", size=13, bold=True, color=PRIMARY)
add_bullets(s, Inches(0.5), Inches(1.6), Inches(6.1), Inches(5.0), [
    ("PyPTO GitHub：github.com/hw-native-sys/pypto", 0),
    ("产品规划：Product_Planning/PyPTO3.0_Toolkit_产品功能规划.md", 0),
    ("仓库洞察：Insight/hw-native-sys-仓库功能洞察报告.md", 0),
    ("TileLang-Ascend：github.com/tile-ai/tilelang-ascend", 0),
    ("Triton-Ascend：github.com/triton-lang/triton-ascend", 0),
    ("Triton 官方(含 Gluon)：triton-lang.org", 0),
    ("TileLang 官方：tilelang.com", 0),
    ("MindStudio：hiascend.com 对应文档页", 0),
], size=11, gap=5)
add_text(s, Inches(6.9), Inches(1.2), Inches(6), Inches(0.3), "战略标杆 / 邻近参照", size=13, bold=True, color=PRIMARY)
add_bullets(s, Inches(6.9), Inches(1.6), Inches(5.9), Inches(5.0), [
    ("cuTile Python：docs.nvidia.com/cuda/cutile-python", 0),
    ("Nsight Compute：docs.nvidia.com/nsight-compute", 0),
    ("CUTLASS / CuTe DSL：docs.nvidia.com/cutlass", 0),
    ("TVM TensorIR / MetaSchedule：tvm.apache.org", 0),
    ("IREE：iree.dev", 0),
    ("vLLM：docs.vllm.ai", 0),
    ("SGLang：docs.sglang.ai", 0),
    ("TensorRT-LLM：nvidia.github.io/TensorRT-LLM", 0),
], size=11, gap=5)

# =========================================================
# Slide 17 — 结尾 / 下一步
# =========================================================
s = new_slide()
add_rect(s, 0, 0, SW, SH, PRIMARY)
add_text(s, Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.8),
         "下一步", size=32, bold=True, color=WHITE)
add_bullets(s, Inches(0.8), Inches(2.5), Inches(11.7), Inches(3.5), [
    ("审核通过 §15 的待确认项后，按“开发者旅程九阶段 + 逐家画像”填充正式内容。", 0, WHITE, False),
    ("补充实测：PyPTO 与 TileLang-Ascend / Triton-Ascend 同机同配置跑定量场景（§13）。", 0, WHITE, False),
    ("可视化：竞品×阶段矩阵（§15）与 DevUX 维度雷达图落地为图表。", 0, WHITE, False),
    ("本 PPT 所有内容均来自项目内已核实文档，未编造数据；规划/推断项已标注。", 0, RGBColor(0xBD,0xD3,0xEE), False),
], size=15, gap=12)
add_text(s, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.4),
         "Competitive_Analysis · 2026-08-11", size=12, color=RGBColor(0x9D,0xB8,0xD8))

# ---------- 保存 ----------
out = "D:/project/PyPTO3/Competitive_Analysis/PyPTO_开发者体验竞品分析.pptx"
prs.save(out)
print("SAVED:", out, "slides:", len(prs.slides._sldIdLst))
